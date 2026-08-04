"""PPTX export helper — reads JSON from stdin, writes .pptx via python-pptx.

Called from workflow nodes via command_exec:
  echo '<json>' | python -m twinkle.agentserver.workflow.ppt_export

Expected JSON shape:
  {"output_path": "output/ppt-xxx/主题.pptx", "topic": "...", "pages": [
      {"title": "...", "body": "...", "page_type": "cover|data|ending"},
      ...
  ]}

Layout rules (Phase 11b-1):
  - cover: 居中标题(Pt44 bold) + 副标题(Pt20) + 日期
  - ending: 居中感谢语(Pt40 bold) + 副文本(Pt18)
  - data: 顶部标题栏(Pt32 bold) + 正文列表(Pt18)
  - Slide size: 13.333" x 7.5" (widescreen 16:9)
"""
from __future__ import annotations

import json
import sys
from datetime import date

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _add_textbox(slide, left, top, width, height):
    """Add a textbox and return its text_frame."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def _fill_first_paragraph(tf, text, size=18, bold=False, alignment=PP_ALIGN.LEFT):
    """Fill the first (default) paragraph in a text_frame."""
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.alignment = alignment
    return p


def generate_pptx(output_path: str, topic: str, pages: list[dict]) -> str:
    """Generate .pptx from page data. Returns output_path on success."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for page in pages:
        page_type = page.get("page_type", "data")
        title = page.get("title", "")
        body = page.get("body", "")

        if page_type == "cover":
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            # Title
            tf = _add_textbox(slide, 1.0, 2.5, 11.333, 1.5)
            _fill_first_paragraph(tf, title, size=44, bold=True, alignment=PP_ALIGN.CENTER)
            # Subtitle
            tf2 = _add_textbox(slide, 1.0, 4.2, 11.333, 1.0)
            _fill_first_paragraph(tf2, body, size=20, alignment=PP_ALIGN.CENTER)
            # Date
            tf3 = _add_textbox(slide, 1.0, 5.5, 11.333, 0.5)
            _fill_first_paragraph(tf3, str(date.today()), size=14, alignment=PP_ALIGN.CENTER)

        elif page_type == "ending":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tf = _add_textbox(slide, 1.0, 2.5, 11.333, 2.0)
            _fill_first_paragraph(tf, title, size=40, bold=True, alignment=PP_ALIGN.CENTER)
            p2 = tf.add_paragraph()
            p2.text = body
            p2.font.size = Pt(18)
            p2.alignment = PP_ALIGN.CENTER

        else:  # data / default
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # Title bar at top
            tf = _add_textbox(slide, 0.5, 0.3, 12.333, 0.8)
            _fill_first_paragraph(tf, title, size=32, bold=True)
            # Body content — each line as a paragraph
            tf2 = _add_textbox(slide, 0.8, 1.5, 11.533, 5.5)
            lines = body.split("\n") if body else [""]
            for i, line in enumerate(lines):
                line = line.strip()
                if not line:
                    continue
                if i == 0:
                    _fill_first_paragraph(tf2, line, size=18)
                else:
                    p = tf2.add_paragraph()
                    p.text = line
                    p.font.size = Pt(18)
                    p.space_after = Pt(8)

    prs.save(output_path)
    return output_path


def main():
    """Entry point: read JSON from stdin, generate PPTX."""
    raw = sys.stdin.read()
    try:
        data = json.loads(raw)
    except json.JSONDecodeError as e:
        print(f"ERROR: invalid JSON: {e}", file=sys.stderr)
        sys.exit(1)

    output_path = data.get("output_path", "")
    if not output_path:
        print("ERROR: output_path is required", file=sys.stderr)
        sys.exit(1)

    topic = data.get("topic", "演示文稿")
    pages = data.get("pages", [])

    try:
        path = generate_pptx(output_path, topic, pages)
        print(f"PPTX saved to {path}")
    except Exception as e:
        print(f"ERROR: pptx generation failed: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
